"""
HTML 프레젠테이션 → PPTX 변환 스크립트

presentation.html의 각 .slide 요소를 헤드리스 브라우저로 렌더링하여
스크린샷을 찍고, 그 이미지를 PPTX 슬라이드에 삽입합니다.
CSS 디자인이 100% 보존됩니다.

사용법:
    python html_to_pptx.py
    python html_to_pptx.py --input my_presentation.html --output result.pptx
"""

import argparse
import os
import shutil
import sys
from pathlib import Path

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Emu


def html_to_pptx(html_path: str, output_path: str):
    html_path = str(Path(html_path).resolve())
    html_url = Path(html_path).as_uri()

    # 임시 폴더에 슬라이드 이미지 저장
    tmp_dir = os.path.join(os.path.dirname(html_path), "_slide_images")
    os.makedirs(tmp_dir, exist_ok=True)

    try:
        with sync_playwright() as p:
            browser = p.chromium.launch(channel="msedge")
            page = browser.new_page(viewport={"width": 1456, "height": 816})
            page.goto(html_url, wait_until="networkidle")

            # 폰트 로딩 대기
            page.wait_for_timeout(2000)

            slides = page.query_selector_all(".slide")
            print(f"슬라이드 {len(slides)}개 발견")

            image_paths = []
            for i, slide in enumerate(slides):
                img_path = os.path.join(tmp_dir, f"slide_{i:03d}.png")
                slide.screenshot(path=img_path)
                image_paths.append(img_path)
                print(f"  슬라이드 {i+1}/{len(slides)} 캡처 완료")

            browser.close()

        # PPTX 생성 (16:9 비율, 1456x816 기준)
        prs = Presentation()
        prs.slide_width = Emu(13716000)   # ~38.1cm
        prs.slide_height = Emu(7772400)   # ~21.59cm (잘못된 비율 보정)

        # HTML의 1456x816 비율에 맞춤
        aspect = 816 / 1456
        actual_aspect = prs.slide_height / prs.slide_width

        for img_path in image_paths:
            slide_layout = prs.slide_layouts[6]  # 빈 슬라이드
            slide = prs.slides.add_slide(slide_layout)

            # 슬라이드 전체를 채우도록 이미지 배치
            slide.shapes.add_picture(
                img_path,
                left=Emu(0),
                top=Emu(0),
                width=prs.slide_width,
                height=prs.slide_height,
            )

        prs.save(output_path)
        print(f"\nPPTX 저장 완료: {output_path}")
        print(f"총 {len(image_paths)}개 슬라이드")

    finally:
        # 임시 이미지 정리
        if os.path.exists(tmp_dir):
            shutil.rmtree(tmp_dir)
            print("임시 파일 정리 완료")


def main():
    parser = argparse.ArgumentParser(description="HTML 슬라이드를 PPTX로 변환")
    parser.add_argument("--input", "-i", default="presentation.html",
                        help="입력 HTML 파일 (기본: presentation.html)")
    parser.add_argument("--output", "-o", default="presentation.pptx",
                        help="출력 PPTX 파일 (기본: presentation.pptx)")
    args = parser.parse_args()

    if not os.path.exists(args.input):
        print(f"오류: 입력 파일을 찾을 수 없습니다: {args.input}")
        sys.exit(1)

    html_to_pptx(args.input, args.output)


if __name__ == "__main__":
    main()
